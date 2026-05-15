"""Lager Presedensdiagram (pptx) basert på avhengighetene i WBS."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

PATH = r"C:\Users\Gusta\OneDrive\Documents\Privat\Studier\Bachelore i Logistikk\LOG565 - Prosjektledelse\02 - Planlegging\Presedensdiagram - Nye Hædda barneskole.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color="000000", align="center", fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor.from_string(fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 2)
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box

def add_box(slide, left, top, width, height, text, fill_color, font_color="FFFFFF", size=10, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    shp.line.color.rgb = RGBColor.from_string("404040")
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 2
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(font_color)
    return shp

def connect(slide, src, dst):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        src.left + src.width, src.top + src.height // 2,
        dst.left, dst.top + dst.height // 2)
    conn.line.color.rgb = RGBColor.from_string("404040")
    conn.line.width = Pt(1.25)
    # Pil
    line = conn.line
    return conn

# === Slide 1: Tittel ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.2),
            "Presedensdiagram", size=44, bold=True, color="1F4E79")
add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.33), Inches(0.8),
            "Nye Hædda barneskole", size=28, color="2E75B6")
add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.5),
            "Avhengigheter mellom hovedleveranser (Activity-on-Node)", size=14, color="595959")
add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.5),
            "Versjon 1.0 — 04.05.2026", size=12, color="A6A6A6")

# === Slide 2: Hovedflyt ===
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6),
            "Hovedflyt — sekvens på nivå 1", size=22, bold=True, color="1F4E79", align="left")
add_textbox(slide, Inches(0.3), Inches(0.75), Inches(12.7), Inches(0.4),
            "Pilene viser FS-avhengigheter (Finish-to-Start). Tid og kostnad fra Bårds estimater settes inn i Gantt.",
            size=11, color="595959", align="left")

# 8 hovedgrener på rad — vis sekvensen
y = Inches(3.0)
boxes = []
positions = [
    ("1 Prosjektledelse",        "1F4E79", Inches(0.3),  y - Inches(2.0)),  # over (kontinuerlig)
    ("2 Prosjektering",          "2E75B6", Inches(0.3),  y),
    ("3 Riving/grunn",           "548235", Inches(2.1),  y),
    ("4 Bygg",                   "C65911", Inches(3.9),  y),
    ("5 Tekniske anlegg",        "BF8F00", Inches(5.7),  y),
    ("7 Inventar (FF&E)",        "7030A0", Inches(7.5),  y),
    ("8 Overtakelse",            "843C0C", Inches(9.3),  y),
    ("6 Utomhus",                "385723", Inches(11.1), y),
]
for navn, color, left, top in positions:
    b = add_box(slide, left, top, Inches(1.7), Inches(0.7), navn, color, size=11)
    boxes.append(b)

# Forbindelseslinjer (FS):
# 2 → 3 → 4 → 5 → 7 → 8
# 6 starter også etter 3 (parallell med 4 og 5)
seq = [(1, 2), (2, 3), (3, 4), (4, 5), (5, 6), (2, 7)]  # indeks
# (boxes-rekkefølge: 0=1, 1=2, 2=3, 3=4, 4=5, 5=7, 6=8, 7=6)
# Avh: 2→3, 3→4, 4→5, 5→7, 7→8, 3→6 (utomhus parallell)
edges = [(1, 2), (2, 3), (3, 4), (4, 5), (5, 6), (2, 7)]
for a, b in edges:
    connect(slide, boxes[a], boxes[b])

# 1 (Prosjektledelse) ligger over og varer hele veien — vis stiplet pil over
add_textbox(slide, Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.4),
            "← Prosjektledelse (1.0) løper hele tiden parallelt →",
            size=11, color="404040", align="center", fill="DEEBF7")

# Slide 3: Hovedavhengigheter med detaljer for kritisk linje
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6),
            "Kritisk linje — fra prosjektering til overtakelse", size=22, bold=True, color="1F4E79", align="left")
add_textbox(slide, Inches(0.3), Inches(0.75), Inches(12.7), Inches(0.4),
            "Forventet kritisk linje: 2.1 → 2.2 → 3.3 → 4.1 → 4.2-4.5 (parallell) → 5.x → 7.x → 8.1 → 8.2",
            size=11, color="595959", align="left")

# Lay out kritisk linje
crit = [
    ("2.1\nDetalj-\nprosjektering",  "2E75B6"),
    ("2.2\nOff. godkjenninger",      "2E75B6"),
    ("3.3\nGrunnarbeid",             "548235"),
    ("4.1\nRåbygg",                  "C65911"),
    ("4.2–4.5\nInnvendig + gymsal",  "C65911"),
    ("5.x\nTekniske anlegg",         "BF8F00"),
    ("7.x\nInventar",                "7030A0"),
    ("8.1\nTesting / prøvedrift",    "843C0C"),
    ("8.2\nFerdigbefaring (BP3)",    "843C0C"),
]
y = Inches(2.5)
boxes2 = []
for i, (navn, color) in enumerate(crit):
    left = Inches(0.3 + 1.45 * i)
    b = add_box(slide, left, y, Inches(1.3), Inches(1.2), navn, color, size=10)
    boxes2.append(b)
for i in range(len(boxes2) - 1):
    connect(slide, boxes2[i], boxes2[i+1])

add_textbox(slide, Inches(0.3), Inches(4.5), Inches(12.7), Inches(0.5),
            "Parallelle løp: 6 Utomhus (kan kjøres parallelt med 4 og 5 etter 3.3) — 1 Prosjektledelse løper hele tiden.",
            size=11, color="404040", align="left")

# Tabell over alle avhengigheter
add_textbox(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.4),
            "Hovedavhengigheter (FS):", size=13, bold=True, color="1F4E79", align="left")
deps = [
    "3 (Riving og grunn) starter etter 2 (Prosjektering ferdig).",
    "4 (Bygg) starter etter 3 (Grunn ferdig).",
    "5 (Tekniske anlegg) starter etter 4.1 (Råbygg tett).",
    "6 (Utomhus) starter etter 3 (parallell med 4 og 5).",
    "7 (Inventar) starter etter 4 (Bygg ferdig).",
    "8 (Overtakelse) starter etter 5 og 7.",
]
for i, t in enumerate(deps):
    add_textbox(slide, Inches(0.5), Inches(5.7 + 0.3 * i), Inches(12.5), Inches(0.3),
                f"• {t}", size=10, color="000000", align="left")

# === Slide 4-N: detaljerte presedensblokker per hovedgren ===
from build_wbs import W
from collections import defaultdict
hovedgren_navn = {
    "1": ("1 Prosjektledelse", "1F4E79"),
    "2": ("2 Prosjektering", "2E75B6"),
    "3": ("3 Forberedelse/riving", "548235"),
    "4": ("4 Bygningsmessig", "C65911"),
    "5": ("5 Tekniske anlegg", "BF8F00"),
    "6": ("6 Utomhus", "385723"),
    "7": ("7 Inventar (FF&E)", "7030A0"),
    "8": ("8 Overtakelse", "843C0C"),
}

# For hver hovedgren: vis nivå 2-elementer + interne avhengigheter
for hg_id, (hg_navn, hg_color) in hovedgren_navn.items():
    barn2 = [w for w in W if w[1] == 2 and w[0].startswith(hg_id + ".")]
    if not barn2:
        continue
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6),
                f"Presedens — {hg_navn}", size=22, bold=True, color=hg_color, align="left")

    # Plasser nivå 2-bokser i en horisontal rad
    n = len(barn2)
    if n == 0:
        continue
    box_w = min(Inches(2.0), Inches(12.0) / n)
    gap = (Inches(12.7) - box_w * n) / max(1, n - 1) if n > 1 else Inches(0)
    bx = []
    for i, w in enumerate(barn2):
        wid, _, navn, *_ = w
        left = Inches(0.3) + (box_w + gap) * i
        b = add_box(slide, left, Inches(2.5), box_w, Inches(1.0),
                    f"{wid}\n{navn}", hg_color, size=9)
        bx.append((wid, b))

    # Tegn avhengigheter mellom nivå 2-bokser (basert på avh-feltet i WBS)
    avh_map = {w[0]: w[5] for w in barn2}
    for wid, box in bx:
        avh_str = avh_map.get(wid, "")
        if not avh_str:
            continue
        for dep in avh_str.split(";"):
            dep = dep.strip()
            for wid2, box2 in bx:
                if wid2 == dep:
                    connect(slide, box2, box)
                    break

    # Legg på beskrivelse av hovedgren
    add_textbox(slide, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.5),
                f"Innhold: {n} nivå-2-elementer i denne hovedgrenen.",
                size=11, color="595959", align="left")

prs.save(PATH)
print(f"Lagret: {PATH}")
print(f"Antall slides: {len(prs.slides)}")
