"""Lager WBS-diagram (pptx) basert på WBS-strukturen."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PATH = r"C:\Users\Gusta\OneDrive\Documents\Privat\Studier\Bachelore i Logistikk\LOG565 - Prosjektledelse\02 - Planlegging\WBS-diagram - Nye Hædda barneskole.pptx"

# Hovedstruktur (nivå 1 og 2 fra WBS)
hovedgrener = [
    ("1. Prosjektledelse",        "1F4E79", ["1.1 Prosjektstyring", "1.2 Kontraktsoppfølging", "1.3 Interessenthåndtering", "1.4 Risiko- og kvalitetsstyring"]),
    ("2. Prosjektering",          "2E75B6", ["2.1 Detaljprosjektering", "2.2 Offentlige godkjenninger", "2.3 Konkurransegrunnlag"]),
    ("3. Forberedelse/riving",    "548235", ["3.1 Miljøsanering", "3.2 Riving", "3.3 Grunnarbeid"]),
    ("4. Bygningsmessig",         "C65911", ["4.1 Råbygg", "4.2 1. etasje (1–4)", "4.3 2. etasje (5–7)", "4.4 3. etasje (8–10)", "4.5 Gymsal"]),
    ("5. Tekniske anlegg",        "BF8F00", ["5.1 VVS", "5.2 Elektro", "5.3 Heis", "5.4 SD-anlegg", "5.5 IKT/sikkerhet"]),
    ("6. Utomhus",                "385723", ["6.1 Lekearealer", "6.2 Sport og fritid", "6.3 Infrastruktur", "6.4 Grøntanlegg"]),
    ("7. Inventar (FF&E)",        "7030A0", ["7.1 Løst inventar", "7.2 Spesialutstyr", "7.3 AV-løsninger"]),
    ("8. Overtakelse",            "843C0C", ["8.1 Testing", "8.2 Ferdigbefaring", "8.3 FDV-dok.", "8.4 Brukeropplæring", "8.5 Sluttrapport"]),
]

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Slide 1: Tittel
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color="000000", align="center", fill=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor.from_string(fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": 1, "center": 2, "right": 3}.get(align, 2)
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box

add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.2),
            "WBS-diagram", size=44, bold=True, color="1F4E79")
add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.33), Inches(0.8),
            "Nye Hædda barneskole", size=28, color="2E75B6")
add_textbox(slide, Inches(0.5), Inches(5.2), Inches(12.33), Inches(0.5),
            "LOG565 Prosjektledelse 2 — Mappeinnlevering", size=14, color="595959")
add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.33), Inches(0.5),
            "Versjon 1.0 — 04.05.2026", size=12, color="A6A6A6")

# Slide 2: Oversikt — alle hovedgrener på én slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(slide, Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6),
            "WBS — Oversikt nivå 1 og 2", size=24, bold=True, color="1F4E79", align="left")
add_textbox(slide, Inches(0.3), Inches(0.7), Inches(12.7), Inches(0.4),
            "Nye Hædda barneskole — total 116 linjer i 4 nivåer", size=12, color="595959", align="left")

# Rotnode i topp midten
def add_box(slide, left, top, width, height, text, fill_color, font_color="FFFFFF", size=11, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    shp.line.color.rgb = RGBColor.from_string("404040")
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = 2
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(font_color)
    return shp

# Rotnode
root = add_box(slide, Inches(5.5), Inches(1.3), Inches(2.3), Inches(0.6),
               "Nye Hædda barneskole", "0F2D45", size=13)

# 8 hovedgrener i rad
n = len(hovedgrener)
total_w = Inches(12.6)
gren_w = Inches(1.5)
gap = (total_w - gren_w * n) / (n - 1)
left0 = Inches(0.3)
top_h = Inches(2.4)
for i, (navn, color, _) in enumerate(hovedgrener):
    left = left0 + (gren_w + gap) * i
    add_box(slide, left, top_h, gren_w, Inches(0.55), navn, color, size=10)

# Underpunkt under hver
sub_top = Inches(3.1)
for i, (navn, color, sub) in enumerate(hovedgrener):
    left = left0 + (gren_w + gap) * i
    for j, s in enumerate(sub):
        sub_t = sub_top + Inches(0.45 * j)
        add_box(slide, left, sub_t, gren_w, Inches(0.4), s, "F2F2F2", font_color="000000", size=8, bold=False)

# Slide 3-N: detaljer per gren med nivå 3 + 4
from build_wbs import W
from collections import defaultdict
by_parent = defaultdict(list)
for w in W:
    wid = w[0]
    if "." in wid:
        parent = wid.rsplit(".", 1)[0]
        by_parent[parent].append(w)

# For hver hovedgren (nivå 1), lag en slide med dens tre
hovedgren_ids = [w for w in W if w[1] == 1]
for w1 in hovedgren_ids:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    wid1, niv1, navn1, beskr1, ansv1, avh1, krav1 = w1
    add_textbox(slide, Inches(0.3), Inches(0.2), Inches(12.7), Inches(0.6),
                f"WBS {wid1} — {navn1}", size=22, bold=True, color="1F4E79", align="left")
    add_textbox(slide, Inches(0.3), Inches(0.7), Inches(12.7), Inches(0.4),
                f"Ansvarlig: {ansv1}", size=11, color="595959", align="left")

    # Finn nivå-2-barn
    barn2 = [w for w in W if w[0].startswith(wid1 + ".") and w[1] == 2]
    n2 = len(barn2)
    if n2 == 0:
        continue
    col_w = Inches(2.5)
    gap2 = (Inches(12.7) - col_w * n2) / max(1, n2 - 1) if n2 > 1 else Inches(0)
    color1 = next(h[1] for h in hovedgrener if h[0].startswith(wid1 + "."))

    for i2, w2 in enumerate(barn2):
        left = Inches(0.3) + (col_w + gap2) * i2
        wid2, _, navn2, *_ = w2
        add_box(slide, left, Inches(1.4), col_w, Inches(0.55),
                f"{wid2} {navn2}", color1, size=10)
        # Barn nivå 3
        barn3 = [w for w in W if w[0].startswith(wid2 + ".") and w[1] == 3]
        for i3, w3 in enumerate(barn3):
            wid3, _, navn3, *_ = w3
            top_y = Inches(2.05 + 0.5 * i3)
            if top_y > Inches(7.0):
                break
            add_box(slide, left, top_y, col_w, Inches(0.45),
                    f"{wid3} {navn3}", "FFF2CC", font_color="000000", size=8, bold=False)

prs.save(PATH)
print(f"Lagret: {PATH}")
print(f"Antall slides: {len(prs.slides)}")
